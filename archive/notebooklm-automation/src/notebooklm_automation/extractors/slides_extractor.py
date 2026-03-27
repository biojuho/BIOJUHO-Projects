"""Google Slides text extraction — via Google Slides API.

Usage:
    from notebooklm_automation.extractors.slides_extractor import extract_slides_text

    text = await extract_slides_text(presentation_id)
    text = extract_slides_text_from_export(pptx_bytes)
"""

from __future__ import annotations

import io

from loguru import logger as log


async def extract_slides_text(
    presentation_id: str,
    *,
    credentials_json: str = "",
) -> str:
    """Extract text from a Google Slides presentation via its ID.

    Requires Google Slides API credentials (service account or OAuth).

    Args:
        presentation_id: The Google Slides presentation ID (from URL).
        credentials_json: Path to service account JSON. Falls back to
            ``GOOGLE_APPLICATION_CREDENTIALS`` env var.
    """
    try:
        from google.oauth2 import service_account  # type: ignore[import-untyped]
        from googleapiclient.discovery import build  # type: ignore[import-untyped]
    except ImportError:
        log.warning("[Slides] google-api-python-client not installed")
        return ""

    try:
        import os

        creds_path = credentials_json or os.getenv("GOOGLE_APPLICATION_CREDENTIALS", "")
        if not creds_path:
            log.warning("[Slides] No credentials configured")
            return ""

        creds = service_account.Credentials.from_service_account_file(
            creds_path,
            scopes=["https://www.googleapis.com/auth/presentations.readonly"],
        )
        service = build("slides", "v1", credentials=creds)
        presentation = service.presentations().get(presentationId=presentation_id).execute()

        texts: list[str] = []
        for slide in presentation.get("slides", []):
            slide_texts: list[str] = []
            for element in slide.get("pageElements", []):
                shape = element.get("shape", {})
                text_content = shape.get("text", {})
                for text_element in text_content.get("textElements", []):
                    text_run = text_element.get("textRun", {})
                    content = text_run.get("content", "").strip()
                    if content:
                        slide_texts.append(content)
            if slide_texts:
                texts.append("\n".join(slide_texts))

        full_text = "\n\n---\n\n".join(texts)
        log.info("[Slides] Extracted %d chars from %d slides", len(full_text), len(texts))
        return full_text

    except Exception as e:
        log.warning("[Slides] API extraction failed: %s", e)
        return ""


def extract_slides_text_from_export(data: bytes) -> str:
    """Extract text from an exported PPTX file (binary).

    Uses python-pptx to parse the PowerPoint format.
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        log.warning("[Slides] python-pptx not installed — skipping")
        return ""

    try:
        prs = Presentation(io.BytesIO(data))
        slides_text: list[str] = []

        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            texts.append(row_text)
            if texts:
                slides_text.append("\n".join(texts))

        full_text = "\n\n---\n\n".join(slides_text)
        log.info("[Slides] PPTX extracted %d chars from %d slides", len(full_text), len(slides_text))
        return full_text

    except Exception as e:
        log.warning("[Slides] PPTX extraction failed: %s", e)
        return ""
